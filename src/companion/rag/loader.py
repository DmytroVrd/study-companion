import logging
import re
import zipfile
from io import BytesIO
from pathlib import Path

import httpx
import pdfplumber
from docx import Document as DocxDocument
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from PIL import Image, ImageOps
from pptx import Presentation
from pytesseract import TesseractError, image_to_string

from companion.config import settings
from companion.llm.providers import LLMProviderError, gemini_ocr_image

SUPPORTED_EXTENSIONS = {".pdf", ".txt", ".docx", ".pptx"}
logger = logging.getLogger(__name__)

_IMAGE_MIME_TYPES = {
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".webp": "image/webp",
}


def _clean_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def load_pdf(path: str | Path) -> list[Document]:
    docs: list[Document] = []
    source = str(path)
    with pdfplumber.open(source) as pdf:
        for index, page in enumerate(pdf.pages):
            simple_text = page.extract_text(x_tolerance=1, y_tolerance=3) or ""
            layout_text = page.extract_text(layout=True) or ""
            text = layout_text if len(layout_text) > len(simple_text) else simple_text
            text = _clean_text(text)
            if text:
                docs.append(
                    Document(
                        page_content=f"Page {index + 1}\n{text}",
                        metadata={
                            "page": index + 1,
                            "source": source,
                            "source_type": "pdf",
                            "loc_label": f"page {index + 1}",
                        },
                    )
                )
    return docs


def load_txt(path: str | Path) -> list[Document]:
    source = str(path)
    raw = Path(path).read_text(encoding="utf-8", errors="ignore")
    text = _clean_text(raw)
    if not text:
        return []
    return [
        Document(
            page_content=text,
            metadata={"page": 1, "source": source, "source_type": "txt", "loc_label": "text"},
        )
    ]


def _local_ocr_image(image_bytes: bytes) -> str:
    if not settings.ENABLE_LOCAL_OCR:
        return ""

    try:
        with Image.open(BytesIO(image_bytes)) as image:
            prepared = ImageOps.grayscale(image)
            prepared = ImageOps.autocontrast(prepared)
            prepared = prepared.resize((prepared.width * 2, prepared.height * 2))
            texts = [
                image_to_string(prepared, lang="ukr+eng", config="--psm 6"),
                image_to_string(prepared, lang="ukr+eng", config="--psm 11"),
            ]
    except (OSError, TesseractError, RuntimeError) as exc:
        logger.info("Local OCR skipped: %s", exc)
        return ""

    lines = []
    seen = set()
    for text in texts:
        for line in _clean_text(text).splitlines():
            normalized = re.sub(r"\W+", "", line.lower())
            if len(normalized) < 3 or normalized in seen:
                continue
            seen.add(normalized)
            lines.append(line)
    return _clean_text("\n".join(lines))


def _gemini_ocr_image(image_bytes: bytes, mime_type: str, path: str | Path) -> str:
    if not settings.ENABLE_GEMINI_OCR or not settings.GEMINI_API_KEY:
        return ""

    try:
        return _clean_text(gemini_ocr_image(image_bytes, mime_type))
    except httpx.HTTPStatusError as exc:
        if exc.response.status_code == 429:
            logger.info("Gemini OCR rate-limited for %s.", path)
            return ""
        raise


def _ocr_office_images(path: str | Path, media_prefix: str) -> list[tuple[int, str, str]]:
    if not settings.ENABLE_LOCAL_OCR and not settings.ENABLE_GEMINI_OCR:
        return []

    extracted: list[tuple[int, str, str]] = []
    try:
        with zipfile.ZipFile(path) as archive:
            image_names = [
                name
                for name in archive.namelist()
                if name.startswith(media_prefix)
                and Path(name).suffix.lower() in _IMAGE_MIME_TYPES
            ]
            for index, name in enumerate(image_names[: settings.OCR_MAX_IMAGES_PER_FILE], start=1):
                suffix = Path(name).suffix.lower()
                image_bytes = archive.read(name)
                image_text = _local_ocr_image(image_bytes)
                if len(image_text) < 80:
                    image_text = _gemini_ocr_image(image_bytes, _IMAGE_MIME_TYPES[suffix], path)
                if image_text:
                    extracted.append((index, Path(name).name, image_text))
    except (zipfile.BadZipFile, OSError, LLMProviderError, ValueError, httpx.HTTPError) as exc:
        logger.info("Image OCR skipped for %s: %s", path, exc)
    return extracted


def load_docx(path: str | Path) -> list[Document]:
    source = str(path)
    document = DocxDocument(source)
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    ocr_parts = _ocr_office_images(path, "word/media/")
    parts.extend(f"Image OCR {index} ({name})\n{text}" for index, name, text in ocr_parts)

    text = _clean_text("\n".join(parts))
    docs: list[Document] = []
    if text:
        docs.append(
            Document(
                page_content=text,
                metadata={
                    "page": 1,
                    "source": source,
                    "source_type": "docx",
                    "loc_label": "document",
                },
            )
        )
    for index, name, image_text in ocr_parts:
        docs.append(
            Document(
                page_content=f"Image OCR {index} ({name})\n{image_text}",
                metadata={
                    "page": 1,
                    "source": source,
                    "source_type": "docx_image",
                    "image": index,
                    "loc_label": f"image OCR {index}",
                },
            )
        )
    return docs


def load_pptx(path: str | Path) -> list[Document]:
    source = str(path)
    presentation = Presentation(source)
    docs: list[Document] = []
    ocr_parts = _ocr_office_images(path, "ppt/media/")
    for index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
        text = _clean_text("\n".join(parts))
        if text:
            docs.append(
                Document(
                    page_content=f"Slide {index}\n{text}",
                    metadata={
                        "slide": index,
                        "source": source,
                        "source_type": "pptx",
                        "loc_label": f"slide {index}",
                    },
                )
            )
    for index, name, image_text in ocr_parts:
        docs.append(
            Document(
                page_content=f"Image OCR {index} ({name})\n{image_text}",
                metadata={
                    "source": source,
                    "source_type": "pptx_image",
                    "image": index,
                    "loc_label": f"image OCR {index}",
                },
            )
        )
    return docs


def load_document(path: str | Path) -> list[Document]:
    extension = Path(path).suffix.lower()
    if extension == ".pdf":
        return load_pdf(path)
    if extension == ".txt":
        return load_txt(path)
    if extension == ".docx":
        return load_docx(path)
    if extension == ".pptx":
        return load_pptx(path)
    supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
    raise ValueError(f"Unsupported file type '{extension}'. Supported types: {supported}.")


def split_documents(docs: list[Document]) -> list[Document]:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1200,
        chunk_overlap=180,
        separators=["\n\n", "\n", ".", " "],
    )
    return splitter.split_documents(docs)
