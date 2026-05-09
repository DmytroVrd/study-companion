from pathlib import Path

from docx import Document as DocxDocument
from langchain_core.documents import Document
from pptx import Presentation
from reportlab.pdfgen import canvas


def _write_pdf(path: Path, text: str) -> None:
    pdf = canvas.Canvas(str(path))
    pdf.drawString(72, 720, text)
    pdf.save()


def test_load_pdf_and_split_documents(tmp_path) -> None:
    from companion.rag.loader import load_pdf, split_documents

    pdf_path = tmp_path / "notes.pdf"
    _write_pdf(pdf_path, "Gradient descent updates model parameters using gradients.")

    docs = load_pdf(pdf_path)
    chunks = split_documents(docs)

    assert len(docs) == 1
    assert "Gradient descent" in docs[0].page_content
    assert chunks
    assert "gradients" in " ".join(chunk.page_content for chunk in chunks)


def test_load_txt_docx_and_pptx(tmp_path) -> None:
    from companion.rag.loader import load_document

    txt_path = tmp_path / "notes.txt"
    txt_path.write_text("Psychology helps people understand emotions.", encoding="utf-8")

    docx_path = tmp_path / "notes.docx"
    docx = DocxDocument()
    docx.add_paragraph("Machine learning models learn patterns from data.")
    docx.save(docx_path)

    pptx_path = tmp_path / "slides.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Photosynthesis converts light into energy."
    presentation.save(pptx_path)

    txt_docs = load_document(txt_path)
    docx_docs = load_document(docx_path)
    pptx_docs = load_document(pptx_path)

    assert "emotions" in txt_docs[0].page_content
    assert "patterns" in docx_docs[0].page_content
    assert "Photosynthesis" in pptx_docs[0].page_content
    assert pptx_docs[0].metadata["loc_label"] == "slide 1"


def test_vectorstore_add_and_query_with_hash_embeddings(tmp_path, monkeypatch) -> None:
    import companion.rag.vectorstore as vectorstore

    monkeypatch.setattr(vectorstore.settings, "CHROMA_PERSIST_DIR", str(tmp_path / "chroma"))
    monkeypatch.setattr(vectorstore.settings, "OPENROUTER_API_KEY", "")
    monkeypatch.setattr(vectorstore.settings, "OPENROUTER_EMBEDDING_MODEL", "")
    vectorstore.get_chroma_client.cache_clear()
    vectorstore.get_embeddings.cache_clear()

    collection = "test_collection"
    docs = [
        Document(page_content="Neural networks learn weights with backpropagation."),
        Document(page_content="Photosynthesis happens in plants."),
    ]

    added = vectorstore.add_documents(collection, docs)
    results = vectorstore.query(collection, "How do networks learn weights?", k=1)

    assert added == 2
    assert len(results) == 1
    assert "weights" in results[0].page_content


def test_vectorstore_reranks_exact_keyword_matches(tmp_path, monkeypatch) -> None:
    import companion.rag.vectorstore as vectorstore

    monkeypatch.setattr(vectorstore.settings, "CHROMA_PERSIST_DIR", str(tmp_path / "chroma"))
    monkeypatch.setattr(vectorstore.settings, "OPENROUTER_API_KEY", "")
    monkeypatch.setattr(vectorstore.settings, "OPENROUTER_EMBEDDING_MODEL", "")
    vectorstore.get_chroma_client.cache_clear()
    vectorstore.get_embeddings.cache_clear()

    collection = "ranking_collection"
    docs = [
        Document(page_content="General psychology includes many useful ideas."),
        Document(
            page_content=(
                "Psychological knowledge is understanding how people think, "
                "feel, and interact with each other."
            )
        ),
    ]

    vectorstore.add_documents(collection, docs)
    results = vectorstore.query(collection, "what is psychological knowledge?", k=1)

    assert "think, feel, and interact" in results[0].page_content
